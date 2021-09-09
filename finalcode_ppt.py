from wand.display import display , Image     # wand library for link to ImageMagick software
             
import os
                 
os.chdir('E:\python\python_script')


#part -1 
#Resize images

def resizeImages():
    convert('nike_black.png', 'transformed_Nike.png', 'x30')      # resize the nike picture to new ratio
    for i in range(1,6) :
        i = str(i)
        imageName = 'image'+i+'.jpg'
        transformedImageName = 'transformedImage_'+i+'.png'
        size = 'x300'
        convert(imageName, transformedImageName, size)            # Calling the convert method to edit

def convert(imageName, transformedImageName, size):
    with Image(filename = imageName) as image :
        image.transform(resize = size)
        image.save(filename=transformedImageName)
        display(image)

 

#part -2 
#Creating the slides  
      
from pptx import Presentation                     # pptx library for create persentation
from pptx.util import Inches                      # Inches for editing the image size (aspect Ratio)

def ppt_create():
    ppt = Presentation()
    slide_register = ppt.slide_layouts[1]
    for i in range(1,6):
        slide_create(ppt, slide_register, str(i))           ## passing parameter to slide_create method 
    ppt.save('Presentation.pptx')

def slide_create(ppt, slideRegister, slideNumber):
    slide = ppt.slides.add_slide(slideRegister)
    img = 'transformedImage_'+slideNumber+'.png'
    from_top = Inches(3)
    from_left = Inches(1)
    add_picture = slide.shapes.add_picture(img , from_left,from_top)
    img2 = 'transformed_Nike.png'
    from_top = Inches(3.3)
    from_left = Inches(1.2)
    add_picture = slide.shapes.add_picture(img2 , from_left , from_top)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, This is Slide " + slideNumber
    subtitle.text = "python-pptx was here!"



# Calling the function to resize and create the ppt 
resizeImages()
ppt_create()


